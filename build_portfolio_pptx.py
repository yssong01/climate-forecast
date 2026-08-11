"""
build_portfolio_pptx.py — 포트폴리오 발표 자료 생성 (19슬라이드).

'정리파일/2-1. 개인 프로젝트 climate-forecast_송영상_ver1-1.pptx' 의 서식을
그대로 읽어 확장한다. 템플릿에서 추출한 사양:

  슬라이드 13.33 x 7.5in (16:9), Blank 레이아웃
  상단 바 #2D2D2D / 제목 맑은고딕 22pt Bold / 부제 11pt #5A5A5A
  소제목 15pt Bold #4A148C / 본문 12.5pt #262626
  결론 박스 #F4EFF7 + 좌측 액센트 바 #4A148C
  이미지 자리표시자 우측 33% (8.50, 1.38, 4.45 x 4.95in)

문체는 학술 서술체로 통일한다 — 구어적 표현·수사의문문·과장을 쓰지 않고,
측정된 사실과 그 근거만 간결하게 기술한다.

**수치는 체크포인트에서 직접 읽는다.** 손으로 옮겨 적으면 재학습 때마다
낡는다 — 실제로 README 결과 요약이 초기 30일 수집본 기준에 멈춰 있다가
2026-08-11 점검에서 발견됐다. 같은 실수를 반복하지 않기 위해 load_metrics()
가 체크포인트를 열어 MAE·F1·기준선·게이트 배분을 가져온다.

실행: python build_portfolio_pptx.py     (torch 필요 — GPU 컨테이너에서 실행)
"""
import os

import torch
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

TEMPLATE = "정리파일/2-1. 개인 프로젝트 climate-forecast_송영상_ver1-1.pptx"
OUT      = "정리파일/2-1. 개인 프로젝트 climate-forecast_송영상_ver1-2.pptx"

CKPT_TRI    = "./checkpoints/numerical_trichef.pt"
CKPT_ZONLY  = "./checkpoints/z_only_current_dataset.pt"

FONT     = "맑은 고딕"
INK      = RGBColor(0x2D, 0x2D, 0x2D)
SUB      = RGBColor(0x5A, 0x5A, 0x5A)
BODY     = RGBColor(0x26, 0x26, 0x26)
ACCENT   = RGBColor(0x4A, 0x14, 0x8C)
RULE     = RGBColor(0xE0, 0xE0, 0xE0)
KBOX_BG  = RGBColor(0xF4, 0xEF, 0xF7)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
TBL_HEAD = RGBColor(0xEC, 0xE4, 0xF2)
WARN     = RGBColor(0xB7, 0x3E, 0x0A)

BODY_L        = 0.85
BODY_T        = 1.45
BODY_W_FULL   = 11.65
BODY_W_NARROW = 7.45          # 우측 이미지가 붙는 슬라이드의 본문 폭
IMG_BOX       = (8.50, 1.38, 4.45, 4.95)   # 슬라이드 폭의 약 33%
KBOX_TOP      = 6.42


# ── 체크포인트에서 실측값 로드 ────────────────────────────────────────

def load_metrics():
    """3축·Z단독 체크포인트에서 발표에 쓸 수치를 읽는다."""
    m = {}
    tri = torch.load(CKPT_TRI, map_location="cpu", weights_only=True)
    m["tri_temp"]   = tri["val_temp_mae"]
    m["tri_precip"] = tri["val_precip_mae"]
    m["epoch"]      = tri.get("epoch")
    m["embed_dim"]  = tri.get("embed_dim", 64)
    d = tri.get("diagnostics") or {}
    m["w_re"], m["w_im"], m["w_z"] = d.get("w_re"), d.get("w_im"), d.get("w_z")
    m["cos_re_z"]  = d.get("cos_re_z_post")
    m["cos_re_im"] = d.get("cos_re_im_post")
    m["cos_im_z"]  = d.get("cos_im_z_post")
    m["extreme"]   = tri.get("extreme_metrics") or {}
    # 기준선 — train.py 수정 이후 체크포인트부터 저장된다
    m["naive_temp"]   = tri.get("val_temp_naive_mae")
    m["naive_precip"] = tri.get("val_precip_naive_mae")
    m["tri_mean"] = tri.get("mean")

    if os.path.exists(CKPT_ZONLY):
        z = torch.load(CKPT_ZONLY, map_location="cpu", weights_only=True)
        m["z_temp"]   = z["val_temp_mae"]
        m["z_precip"] = z["val_precip_mae"]
        m["z_extreme"] = z.get("extreme_metrics") or {}
        m["z_mean"] = z.get("mean")
        # 같은 데이터셋에서 학습됐는지 — 다르면 비교 자체가 성립하지 않는다
        m["same_dataset"] = (m["tri_mean"] == m["z_mean"])
    else:
        m["z_temp"] = m["z_precip"] = None
        m["same_dataset"] = None
    return m


def pct(new, old):
    """old 대비 new 의 변화율(%). 오차 지표이므로 음수가 개선이다."""
    if not old or not new:
        return "—"
    return f"{(new - old) / old * 100:+.1f}%"


def fnum(v, fmt="{:.3f}", dash="—"):
    return fmt.format(v) if isinstance(v, (int, float)) else dash


# ── 저수준 도우미 ────────────────────────────────────────────────────

def _txbox(slide, l, t, w, h):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return box, tf


def _run(para, text, size, bold=False, color=BODY, italic=False):
    r = para.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return r


def _rect(slide, l, t, w, h, fill=None, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


def new_slide(prs, title, tagline, runhead, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, 13.33, 0.06, fill=INK)
    _, tf = _txbox(slide, 0.5, 0.16, 10.6, 0.5)
    _run(tf.paragraphs[0], title, 21, bold=True, color=INK)
    _, tf = _txbox(slide, 0.64, 0.6, 11.1, 0.3)
    _run(tf.paragraphs[0], tagline, 11, color=SUB)
    _rect(slide, 0.5, 0.88, 12.3, 0.012, fill=RULE)
    _, tf = _txbox(slide, 9.6, 0.2, 3.3, 0.34)
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    _run(tf.paragraphs[0], f"{runhead}   {page} / 19", 10, color=SUB)
    return slide


def subtitle(slide, text, top=1.02, left=0.59, width=12.24):
    _, tf = _txbox(slide, left, top, width, 0.37)
    _run(tf.paragraphs[0], text, 14, bold=True, color=ACCENT)


def kbox(slide, segs):
    _rect(slide, 0.59, KBOX_TOP, 12.24, 0.66, fill=KBOX_BG)
    _rect(slide, 0.59, KBOX_TOP, 0.05, 0.66, fill=ACCENT)
    _, tf = _txbox(slide, 0.78, KBOX_TOP + 0.08, 11.9, 0.54)
    p = tf.paragraphs[0]
    for seg, bold in segs:
        _run(p, seg, 12, bold=bold, color=ACCENT)


def img_placeholder(slide, title, caption, box=None):
    """이미지 삽입 자리 — 제목 + 설명글을 함께 표시해 무엇을 넣을지 명확히 한다."""
    l, t, w, h = box or IMG_BOX
    sh = _rect(slide, l, t, w, h, fill=None, line=ACCENT)
    sh.line.width = Pt(1.5)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, "[ 이미지 삽입 ]\n", 11, bold=True, color=ACCENT)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    _run(p2, title + "\n", 12, bold=True, color=INK)
    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    _run(p3, caption, 9.5, color=SUB)


def bullets(slide, items, left=BODY_L, top=BODY_T, width=BODY_W_FULL,
            size=11.5, bottom=None):
    avail = max(0.4, (bottom or KBOX_TOP - 0.08) - top)
    _, tf = _txbox(slide, left, top, width, avail)
    first = True
    for lvl, segs in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        p.space_after = Pt(4)
        p.line_spacing = 1.2
        mark = "• " if lvl == 0 else "– "
        _run(p, ("" if lvl == 0 else "    ") + mark, size,
             color=ACCENT if lvl == 0 else SUB)
        for text, bold in segs:
            _run(p, text, size, bold=bold, color=BODY)
    return tf


def table(slide, rows, left, top, width, col_w=None, size=9.5,
          head_size=9.5, row_h=0.28):
    n_r, n_c = len(rows), len(rows[0])
    gf = slide.shapes.add_table(n_r, n_c, Inches(left), Inches(top),
                                Inches(width), Inches(row_h * n_r))
    tbl = gf.table
    tbl.first_row = True
    if col_w:
        for i, w in enumerate(col_w):
            tbl.columns[i].width = Inches(w)
    for ri, row in enumerate(rows):
        tbl.rows[ri].height = Inches(row_h)
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.01)
            cell.margin_bottom = Inches(0.01)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = TBL_HEAD if ri == 0 else WHITE
            p = cell.text_frame.paragraphs[0]
            cell.text_frame.word_wrap = True
            segs, buf, bold = [], "", False
            i = 0
            while i < len(cell_text):
                if cell_text[i:i+2] == "**":
                    if buf:
                        segs.append((buf, bold)); buf = ""
                    bold = not bold
                    i += 2
                else:
                    buf += cell_text[i]; i += 1
            if buf:
                segs.append((buf, bold))
            for text, b in segs:
                _run(p, text, head_size if ri == 0 else size,
                     bold=True if ri == 0 else b,
                     color=ACCENT if ri == 0 else BODY)
    return gf


def code_line(slide, text, left, top, width, size=12.5, height=0.38):
    box = _rect(slide, left, top, width, height, fill=KBOX_BG)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Consolas"
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = ACCENT


def stat_row(slide, stats, left, top, width, gap=0.12, h=0.8):
    n = len(stats)
    w = (width - gap * (n - 1)) / n
    for i, (v, k) in enumerate(stats):
        x = left + i * (w + gap)
        _rect(slide, x, top, w, h, fill=KBOX_BG)
        _, tf = _txbox(slide, x + 0.08, top + 0.08, w - 0.16, h - 0.14)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _run(p, v, 17, bold=True, color=ACCENT)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        _run(p2, k, 9, color=SUB)



# ── 본문 ─────────────────────────────────────────────────────────────

def build():
    if not os.path.exists(TEMPLATE):
        raise SystemExit(f"템플릿을 찾을 수 없습니다: {TEMPLATE}")
    M = load_metrics()

    prs = Presentation(TEMPLATE)
    xml_slides = prs.slides._sldIdLst
    for sld in list(xml_slides):
        prs.part.drop_rel(sld.rId)
        xml_slides.remove(sld)

    E = M["extreme"]
    ZE = M.get("z_extreme") or {}
    f1 = lambda k: fnum(E.get(k, {}).get("f1"))
    zf1 = lambda k: fnum(ZE.get(k, {}).get("f1"))
    pr = lambda k: fnum(E.get(k, {}).get("precision"), "{:.1%}")
    rc = lambda k: fnum(E.get(k, {}).get("recall"), "{:.1%}")
    npos = lambda k: f"{E.get(k, {}).get('n_pos', 0):,}건"

    W = BODY_W_NARROW      # 우측 이미지가 있는 슬라이드의 본문 폭
    F = BODY_W_FULL

    # ══════════════════════════════════ 1. 표지 · 종합 요약
    s = new_slide(prs, "환경 예보 (climate-forecast)",
                  "지상 관측 12개 지점을 기반으로 +6시간 뒤 기온·강수 및 극한기상 확률을 산출하는 3축 멀티모달 파이프라인",
                  "개요", 1)
    stat_row(s, [
        (f"{M['tri_temp']:.2f}°C", "기온 MAE"),
        (f"{M['tri_precip']:.3f}mm", "강수 MAE"),
        (pct(M["tri_temp"], M["z_temp"]), "기온 개선\n(Z축 단독 대비)"),
        (pct(M["tri_precip"], M["z_precip"]), "강수 개선\n(Z축 단독 대비)"),
        (f1("heatwave"), "폭염 F1"),
        ("332,195", "학습 관측 건수"),
    ], BODY_L, 1.28, F, h=0.78)

    subtitle(s, "연구 개요", top=2.22, left=BODY_L, width=7.0)
    table(s, [
        ["구분", "내용"],
        ["문제", "지상 관측망은 공간적으로 성글며, 단일 시점의 관측만으로는 상태의 변화 방향을 판별할 수 없다"],
        ["접근", "Tri-CHEF 3축 융합 수식을 이식하여 주변(공간 보간)·흐름(시간 경향)·현재(수치 관측)를 분리 결합"],
        ["결과", f"Z축 단독 대비 기온 {pct(M['tri_temp'], M['z_temp'])}, 강수 {pct(M['tri_precip'], M['z_precip'])} 개선. "
                 f"극한기상 3종 분류 헤드 추가 및 상시 배포"],
    ], BODY_L, 2.62, W, col_w=[0.85, 6.6], row_h=0.5)

    subtitle(s, "구성", top=4.28, left=BODY_L, width=7.0)
    table(s, [
        ["항목", "내용"],
        ["데이터", "기상자료개방포털 ASOS 시간자료 332,195건 (2023-01~2026-08, 12개 지점)"],
        ["모델", f"3축 Tri-CHEF (embed_dim {M['embed_dim']}), 동적 게이팅, 회귀 2·분류 3 헤드"],
        ["파이프라인", "포털 파일셋 수집 → GPU 학습 → CPU 서빙 → 6시간 주기 자동 갱신"],
        ["환경", "PyTorch 2.6 / Docker / Streamlit Cloud (서빙 메모리 227.9MiB 측정)"],
        ["저장소", "https://github.com/yssong01/climate-forecast"],
    ], BODY_L, 4.68, W, col_w=[1.0, 6.45], row_h=0.28)

    img_placeholder(s, "① 대시보드 상단",
                    "관측 지점명과 현재 기온·강수·습도·기압\n4개 지표가 표시되는 영역",
                    box=(8.50, 2.22, 4.45, 4.1))
    kbox(s, [("논문의 3축 융합 수식을 기후 예측 도메인에 이식하고, 그 효과를 ", False),
             ("동일 데이터·동일 분할 조건에서 정량 검증", True), ("하였다.", False)])

    # ══════════════════════════════════ 2. 문제 정의
    s = new_slide(prs, "2-1. 문제 정의 및 연구 범위",
                  "관측 데이터의 구조적 한계와 이에 대응하여 설정한 범위", "문제 정의", 2)
    subtitle(s, "관측 데이터의 세 가지 한계")
    table(s, [
        ["한계", "내용", "설계 대응"],
        ["**공간적 희소성**", "전국 ASOS 주요 지점은 수십 개 수준이며, 지점 사이 구간은 관측이 존재하지 않는다",
         "인접 지점의 동일 시각 관측을 공간 보간하여 Re축으로 구성"],
        ["**시간 정보의 부재**", "동일한 1005hPa라도 상승 중인 경우와 하강 중인 경우의 의미가 다르나, 단일 레코드에는 그 구분이 없다",
         "1·3·6시간 전 대비 변화량을 Im축으로 구성"],
        ["**사건의 희소성**", "폭염·한파·황사는 전체의 수 % 수준이므로, 정확도로 평가하면 '미발생'만 답해도 99%에 도달한다",
         "정확도 대신 정밀도·재현율·F1로 평가"],
    ], BODY_L, 1.42, F, col_w=[2.0, 5.5, 4.15], row_h=0.62)
    subtitle(s, "연구 범위 설정과 근거", top=3.65)
    table(s, [
        ["항목", "선택", "근거"],
        ["예측 시계", "+6시간",
         "+1시간 예측은 퍼시스턴스(T(t+1)≈T(t)) 대비 개선 폭이 +0.7%에 그쳐 모델의 기여를 확인할 수 없었다"],
        ["대상 변수", "기온·강수 및 극한기상 3종",
         "회귀와 이진 분류를 동일 표현에서 분기시켜 헤드 간 상호작용을 관측 가능하게 하였다"],
        ["제외", "가뭄·기후변화 감시", "수주~수개월 누적 지표로, +6시간 시계와 시간 척도가 부합하지 않는다"],
        ["제외", "대설 독립 헤드", "한파와 호우의 논리곱에 해당하여 중복 계산이 발생한다"],
        ["제외", "강풍·우박", "학습 표본이 각각 약 300건·9건으로, 표본 부족은 하이퍼파라미터 조정으로 해결되지 않는다"],
    ], BODY_L, 4.02, F, col_w=[1.35, 2.3, 8.0], row_h=0.38)
    kbox(s, [("예측 시계, 대상 변수, 제외 항목을 모두 ", False),
             ("측정 결과에 근거하여 결정", True), ("하였다.", False)])

    # ══════════════════════════════════ 3. 데이터
    s = new_slide(prs, "2-2. 데이터 — 출처, 규모, 수집 경로의 전환",
                  "API 호출 제한을 확인하고 웹 포털 파일셋으로 수집 경로를 변경한 경위", "데이터", 3)
    subtitle(s, "수집 경로 — 기상자료개방포털 (data.kma.go.kr), 시간자료 2023~2026")
    table(s, [
        ["포털 내 경로", "용도"],
        ["데이터 → 기상관측 → 지상 → **종관기상관측(ASOS)**", "Z·Re·Im 3축의 원천 관측값 (332,195건, 유효 쌍 331,595건)"],
        ["데이터 → 기상관측 → 지상 → **황사관측(PM10)**", "황사 분류 헤드의 정답 라벨 (1시간 평균 150㎍/㎥ 기준)"],
        ["**날씨 이슈별 데이터** → 폭염·한파·황사·태풍", "폭염·한파 분류 헤드의 정답 라벨 (발표된 특보 기록)"],
    ], BODY_L, 1.42, F, col_w=[5.5, 6.15], row_h=0.32)
    stat_row(s, [("332,195", "ASOS 시간자료"), ("331,595", "(t, t+6h) 유효 쌍"),
                 ("12", "관측 지점"), ("3.6년", "2023-01 ~ 2026-08"),
                 ("265,276", "학습셋 (80%)"), ("66,319", "검증셋 (20%)")],
             BODY_L, 2.65, F, h=0.76)
    subtitle(s, "수집 경로 전환 — 문제·대응·결과", top=3.62)
    table(s, [
        ["단계", "내용"],
        ["문제", "apihub API를 통한 대량 수집이 **누적 약 9,800건 시점에서 차단**되었다(동시성 16 조건에서 측정). "
                 "병렬화로 우회되지 않는 요청 수 기준의 제한으로 판단된다"],
        ["대응", "과거 데이터 수집은 웹 포털 파일셋 일괄 다운로드로 전환하고(import_kma_fileset.py), "
                 "API는 실시간 추론 용도로만 사용하도록 역할을 분리하였다"],
        ["결과", "30일분에서 **3.6년분(332,195건)**으로 확장하였으며, 계절 순환이 학습셋에 포함되어 "
                 "시간 기반 분할 검증의 제약도 해소되었다"],
    ], BODY_L, 4.02, F, col_w=[0.9, 10.75], row_h=0.54)
    kbox(s, [("원본 CSV는 ", False), ("EUC-KR 인코딩", True),
             ("이며, 결측은 기본값으로 대체하지 않고 해당 레코드 전체를 제외한다 — 결측이 실측으로 오인되는 것을 방지한다.", False)])

    # ══════════════════════════════════ 4. 아키텍처
    s = new_slide(prs, "2-3. 시스템 아키텍처",
                  "수집·학습·서빙·갱신의 역할 분리와 그 근거", "아키텍처", 4)
    subtitle(s, "4개 경로")
    table(s, [
        ["단계", "구성", "산출물"],
        ["① 수집", "포털 파일셋 일괄 수집(과거) + apihub API(실시간, 출력 1건당 15회)", "historical_data_1y.json (69MB)"],
        ["② 학습", "Docker GPU 컨테이너, PyTorch 2.6, 조기 종료 기준 patience 20", "numerical_trichef.pt (236KB)"],
        ["③ 서빙", "Streamlit + CPU torch, 서빙 메모리 **227.9MiB 측정**", "Streamlit Community Cloud"],
        ["④ 갱신", "GitHub Actions 6시간 주기 → 저장소 커밋 → 자동 재배포", "recent_window.json (0.6MB)"],
    ], BODY_L, 1.42, F, col_w=[0.95, 6.7, 4.0], row_h=0.34)
    subtitle(s, "설계 판단", top=3.15)
    bullets(s, [
        (0, [("학습 산출물과 코드 이미지를 분리하였다. ", True),
             ("학습은 관측 수집과 GPU 연산을 요구하여 이미지 빌드 과정에 포함하기 부적합하므로, 체크포인트는 런타임에 볼륨으로 연결한다.", False)]),
        (0, [("배포 산출물에서 학습 캐시를 제외하였다. ", True),
             ("텍스트 임베딩(460MB)과 원본 이력(69MB)을 제외하고 경량 관측 창(0.6MB)만 사용하여 서빙 메모리를 1.55GiB에서 227.9MiB로 축소하였다.", False)]),
        (0, [("서빙 경로에서 torchvision·sentence-transformers 의존성을 제거하였다. ", True),
             ("현행 체크포인트에서 해당 경로가 호출되지 않음을 코드 수준에서 확인하였다.", False)]),
        (0, [("배포 환경의 파일시스템은 휘발성이다. ", True),
             ("애플리케이션이 데이터를 보존하는 대신, 저장소를 원본으로 두고 Actions가 갱신 결과를 커밋하는 구조로 전환하였다.", False)]),
    ], top=3.55, size=11)
    kbox(s, [("배포 환경의 제약(메모리 한도, 파일시스템 휘발성)을 ", False),
             ("아키텍처 설계 단계에서 반영", True), ("하였다.", False)])

    # ══════════════════════════════════ 5. 선행 연구 검토
    s = new_slide(prs, "2-4. 선행 연구 검토 — 기상 AI 모델의 적용 가능성",
                  "격자 기반 모델의 전제 조건과 본 연구 데이터의 불일치", "선행 연구", 5)
    subtitle(s, "검토 대상과 적용을 유보한 근거")
    table(s, [
        ["모델", "핵심 방법", "적용을 유보한 근거"],
        ["**GraphCast** (DeepMind)\n**Pangu-Weather** (Huawei)\n**FourCastNet** (NVIDIA)\n**ClimaX** (Microsoft)",
         "그래프 신경망, 3D Transformer, Fourier Neural Operator 등으로 전지구 격자 상태를 전파",
         "네 모델 모두 **ERA5 재분석 격자(0.25° 전지구)**를 입력 전제로 한다. 본 연구의 입력은 **12개 지점**이므로 "
         "격자 구조가 존재하지 않으며, 지점을 격자로 보간하여 입력할 경우 관측이 아니라 **보간 과정의 인공물을 학습**하게 된다"],
        ["**NowcastNet**", "이류–확산 물리 손실을 적용한 초단기 강수 예측",
         "이류·확산 항을 본 연구의 IDW 보간장에 적용할 경우, **자체 보간 인공물을 물리 법칙으로 학습**하는 결과가 된다"],
    ], BODY_L, 1.42, F, col_w=[2.5, 3.4, 5.75], row_h=1.05)
    subtitle(s, "대신 반영한 기상학·통계학 원리", top=3.9)
    table(s, [
        ["원리", "내용", "적용 위치"],
        ["역거리가중(IDW) 공간 보간", "지구통계학의 표준 보간 기법. 거리 제곱에 반비례하는 가중(power=2)", "Re축 구성"],
        ["기압 경향 (pressure tendency)", "기압 하강률은 저기압 접근의 선행 지표이며, 기온 하강률은 복사냉각 과정을 반영한다", "Im축 구성"],
        ["퍼시스턴스 기준선", "단기 예측 평가의 표준 기준선 T(t+L)≈T(t)", "성능 판정 기준"],
        ["Hurdle 모델", "영과잉(zero-inflated) 분포를 발생 여부와 발생량으로 분해하는 계량경제학·수문학의 표준 접근", "강수 회귀 헤드"],
    ], BODY_L, 4.3, F, col_w=[2.6, 7.05, 2.0], row_h=0.42)
    kbox(s, [("최신 모델을 배제한 것이 아니라, ", False),
             ("입력 데이터가 해당 모델의 전제 조건을 충족하지 못함을 확인하고 적용을 유보", True), ("하였다.", False)])

    # ══════════════════════════════════ 6. Tri-CHEF 이식 근거
    s = new_slide(prs, "2-5. Tri-CHEF 이식 근거",
                  "검색 도메인에서 제안된 융합 수식을 기후 예측에 적용한 이유", "모델 설계", 6)
    subtitle(s, "원 논문 — Complex-Hermitian Embedding Fusion for Korean Multimodal Retrieval")
    bullets(s, [
        (0, [("Zenodo 프리프린트, DOI 10.5281/zenodo.20034370, CC BY 4.0. ", False),
             ("본 연구와 도메인이 상이한 **한국어 멀티모달 검색** 논문이다.", False)]),
        (0, [("문서·이미지·영상·음원을 통합 검색할 때 ", False),
             ("서로 다른 3개 인코더(SigLIP2·BGE-M3·DINOv2)", True),
             ("의 출력을 단일 점수로 결합하는 문제를 다룬다.", False)]),
        (0, [("통상적인 ", False), ("가중합", True),
             ("은 특정 채널의 점수가 클 경우 해당 채널이 최종 점수를 지배한다. 논문은 대안으로 ", False),
             ("Hermitian-style modulus", True),
             ("(제곱–합–제곱근)를 제안하며, 이 구조에서는 한 축이 다른 축을 상쇄하여 소거하지 못하고 세 축의 근거가 모두 보존된다. "
              "이것이 논문의 ", False), ("식 1(Eq.1)", True), ("이다.", False)]),
    ], top=1.42, width=W, size=11)
    subtitle(s, "이식 시 이점과 함께 이식된 한계", top=3.05, left=0.59, width=7.6)
    table(s, [
        ["Tri-CHEF의 성질", "본 연구에서의 의미"],
        ["축별 근거 보존 (단일 축의 지배 불가)",
         "공간·시간·현재 중 일부가 결측되어도 나머지 축이 유지된다. 관측 실패가 빈번한 현장 데이터에 적합하다"],
        ["축 단위 해석 가능성",
         "출력마다 축별 기여도를 산출할 수 있어, 산출 근거를 화면에 제시할 수 있다"],
        ["이종 모달리티 전제",
         "격자(보간장)·벡터(경향)·표(수치)를 동일 임베딩 공간에 투영하는 구조를 그대로 활용하였다"],
        ["**한계 ①** 헤드 간 표현 경쟁",
         "축과 헤드가 증가할수록 공유 표현을 두고 경쟁이 발생하여 강수 성능이 저하되었다"],
        ["**한계 ②** 위상각의 퇴화",
         "세 인코더가 모두 단위 정규화 출력이므로 θ = atan2(w_im, w_re)로 축소된다. 게이트 가중치의 재표현에 불과하며, 두 값의 차이는 1.19e-07로 측정되었다"],
    ], BODY_L, 3.42, W, col_w=[2.5, 4.95], row_h=0.47)
    img_placeholder(s, "② 모델 구조 탭 — 원 논문 관계",
                    "화면에서 원 논문의 도메인과\n본 연구의 차이를 설명하는 영역",
                    box=(8.50, 3.05, 4.45, 3.28))
    kbox(s, [("이식한 것은 ", False), ("결합 방식(수식 구조)", True),
             ("이며, 세 축에 입력되는 데이터는 본 연구에서 새로 설계하였다.", False)])

    # ══════════════════════════════════ 7. 3축 설계
    s = new_slide(prs, "2-6. 3축 설계 — Re·Im·Z의 정의",
                  "단일 시점 관측에서 유도할 수 없는 정보만을 새 축에 배치한다는 기준", "모델 설계", 7)
    subtitle(s, "축별 정의")
    table(s, [
        ["축", "대응 질문", "입력", "Z축에서 유도 불가능한 근거"],
        ["**Z** 현재", "해당 지점의 현재 상태는 무엇인가",
         "기온·강수·습도·풍속·풍향(sin/cos)·기압·강수형태·시각(sin/cos)·위도·경도 = **12차원**", "기준축"],
        ["**Re** 주변", "인접 지역의 현재 상태는 무엇인가",
         "대상 지점을 **제외한** 11개 지점의 동일 시각 관측을 IDW(power=2)로 ±3°(약 300km) 격자 **32×32×4채널**에 보간",
         "타 지점의 관측값은 해당 지점의 레코드에 포함되지 않는다"],
        ["**Im** 경향", "상태가 어느 방향으로 변화하고 있는가",
         "동일 지점의 **1·3·6시간 전 대비 변화량** × 4변수(기온·기압·습도·풍속) = **12차원**",
         "∂s/∂t 는 단일 시점 s(t)만으로 산출할 수 없다"],
    ], BODY_L, 1.42, W, col_w=[0.95, 1.75, 3.0, 1.75], row_h=0.95)
    subtitle(s, "명명의 유래와 지도 표시의 근거", top=4.5, left=0.59, width=7.6)
    bullets(s, [
        (0, [("Re·Im 명칭은 원 논문이 세 축을 복소수에 대응시킨 데서 유래한다(Real·Imaginary). "
              "복소수에서 두 축이 직교한다는 성질을 축 간 정보 중복을 배제한다는 설계 의도로 차용한 것이며, "
              "**실제 연산은 모두 실수 영역에서 수행된다.**", False)]),
        (0, [("설정 화면의 지도가 지점 위치만을 표시하는 이유는 ", False),
             ("위경도가 실제 모델 입력이기 때문이다", True),
             (". 지점별 기후 특성이 상이하므로(강릉의 동해 영향, 제주의 해양성) 좌표 없이는 단일 모델이 12개 지점을 구분할 수 없다. "
              "또한 Re축의 IDW 가중치가 지점 간 거리로 결정되므로, 지도는 인접 지점의 반영 정도를 설명한다.", False)]),
    ], top=4.9, width=W, size=10, bottom=6.3)
    img_placeholder(s, "③ 설정 화면 — 관측 지점 지도",
                    "12개 지점의 위치와 선택 상태.\n외부 타일 서버 없이 렌더링된다")
    kbox(s, [("Re축에는 ", False), ("대상 지점 자신의 관측값을 포함하지 않는다", True),
             (" — 포함할 경우 Z축과 중복되어 해당 축의 정보량이 소실된다.", False)])

    # ══════════════════════════════════ 8. 융합 수식
    s = new_slide(prs, "2-7. 융합 수식 — Eq.1의 구성",
                  "기호별 정의와 연산 단계. 결과 s는 스칼라가 아니라 길이 64의 벡터이다", "모델 설계", 8)
    code_line(s, "s_i  =  √( (w_Re · Re_i)²  +  (w_Im · Im_i)²  +  (w_Z · Z_i)² )",
              BODY_L, 1.38, W, size=12)
    table(s, [
        ["기호", "정의", "본 학습 결과"],
        ["Re, Im, Z", f"각 축을 인코더에 통과시켜 얻은 **길이 {M['embed_dim']}의 벡터**", f"축당 {M['embed_dim']}개"],
        ["i", "벡터 내 원소의 위치. 연산은 원소별로 수행된다", f"1 ~ {M['embed_dim']}"],
        ["w", "축별 가중치. **Σw = 1** (softmax 정규화)",
         f"Re {fnum(M['w_re'])} / Im {fnum(M['w_im'])} / Z {fnum(M['w_z'])}"],
        ["s", "결합 결과. 6개 예측 헤드가 **공유**하는 입력", f"길이 {M['embed_dim']}"],
    ], BODY_L, 1.86, W, col_w=[1.05, 4.3, 2.1], row_h=0.36)
    subtitle(s, "연산 단계의 의미", top=3.72, left=0.59, width=7.6)
    bullets(s, [
        (0, [("① 곱셈 — ", True), ("신뢰도가 낮게 산출된 축의 기여가 축소된다", False)]),
        (0, [("② 제곱 — ", True), ("부호를 제거하여 축 간 상쇄를 방지한다", False)]),
        (0, [("③ 합산 후 제곱근 — ", True),
             ("직교하는 세 성분의 합성 크기를 산출하는 연산에 해당한다", False)]),
        (0, [("각 축 벡터는 결합 전 단위 길이로 정규화된다. 따라서 결과는 축의 절대 크기가 아니라 "
              "**패턴의 방향과 가중치**로 결정된다.", False)]),
        (0, [("논문 표기 s = √(A² + (αB)² + (φC)²)에서 A·B·C는 Re·Im·Z에 대응하며, "
              "α = w_Im ÷ w_Re, φ = w_Z ÷ w_Re 로 환산된다.", False)]),
    ], top=4.12, width=W, size=10.5, bottom=6.3)
    img_placeholder(s, "④ 모델 구조 탭 — 융합 수식",
                    "수식과 기호별 해설표가\n표시되는 영역")
    kbox(s, [("연산은 벡터의 각 원소마다 독립적으로 수행되며, 결과 s 역시 ", False),
             (f"길이 {M['embed_dim']}의 벡터", True), ("이다.", False)])

    # ══════════════════════════════════ 9. 동적 게이팅
    s = new_slide(prs, "2-8. 동적 게이팅 — 축 가중치의 산출",
                  "고정 상수 대신 입력 조건부 재계산을 채택한 근거", "모델 설계", 9)
    subtitle(s, "구조 및 설계 근거")
    table(s, [
        ["항목", "내용"],
        ["구조", "표준화된 Z축 12차원 → Linear(12→16) → GELU → Linear(16→3) → softmax → (w_Re, w_Im, w_Z)"],
        ["초기화", "균등 초기화(1/3). 논문 기본값 (1, 0.4, 0.2)을 적용할 경우 Re축이 초기부터 우위를 점하며, "
                   "엔트로피 정규화 항이 이 순서를 증폭시켜 축의 유용성이 학습되기 전에 고착되는 현상이 관측되었다"],
        ["softmax 채택 근거", "가중치가 자유로울 경우 (w_Re, w_Im, w_Z)를 k배 축소하고 헤드가 1/k로 보정하여 "
                              "크기 정규화가 무효화된다. Σw=1로 고정하면 축들이 한정된 배분을 공유하므로 가중치 결정이 실질적 선택이 된다"],
        ["문제 배경", "학습 가능한 고정 α·φ는 정보량이 없는 축의 가중치를 오히려 증가시켰다(α 0.4 → 0.665). "
                      "훈련 손실만을 최소화할 경우 무용한 축도 훈련셋 암기에 기여하므로 가중치를 낮출 유인이 없다"],
    ], BODY_L, 1.42, W, col_w=[1.5, 5.95], row_h=0.72)
    subtitle(s, "배분값이 학습된 판단임을 보이는 관측", top=4.68, left=0.59, width=7.6)
    table(s, [
        ["시점", "Re축 입력", "Re 배분", "해석"],
        ["v4", "합성 위성 이미지", "**0.000**", "기여 없음으로 학습"],
        ["v5", "실측 IDW 공간 보간장", "**0.516**", "교체 직후 최대 비중으로 전환"],
        ["현행", "실측 IDW 공간 보간장", fnum(M["w_re"]), "Im축 실측화 이후 3축이 균형"],
    ], BODY_L, 5.08, W, col_w=[0.85, 3.0, 1.1, 2.5], row_h=0.3)
    img_placeholder(s, "⑤ 모델 구조 탭 — 축 배분",
                    "출력마다 재계산되는 축별\n기여도 막대그래프")
    kbox(s, [("게이트 배분은 학습을 통해 결정된 값이며, ", False),
             ("0.000에서 0.516으로의 변화", True),
             ("는 입력 데이터 교체의 효과를 직접적으로 보인다.", False)])

    # ══════════════════════════════════ 10. 적용 효과
    s = new_slide(prs, "2-9. Tri-CHEF 적용 효과 — 3축 대 Z축 단독",
                  "동일 데이터·동일 분할·동일 조기 종료 기준에서 측정한 구조적 효과", "적용 효과", 10)
    subtitle(s, "비교 결과")
    table(s, [
        ["구성", "기온 MAE", "강수 MAE", "폭염 F1", "한파 F1", "황사 F1"],
        ["**Z축 단독** (Re·Im 입력 차단)",
         f"{fnum(M['z_temp'], '{:.4f}')} °C", f"{fnum(M['z_precip'], '{:.4f}')} mm",
         zf1("heatwave"), zf1("coldwave"), zf1("dust")],
        ["**3축 Tri-CHEF** (현행)",
         f"**{M['tri_temp']:.4f} °C**", f"**{M['tri_precip']:.4f} mm**",
         f"**{f1('heatwave')}**", f"**{f1('coldwave')}**", f"**{f1('dust')}**"],
        ["변화",
         f"**{pct(M['tri_temp'], M['z_temp'])}**", f"**{pct(M['tri_precip'], M['z_precip'])}**",
         "+0.090", "**+0.278**", "**+0.382**"],
    ], BODY_L, 1.42, F, col_w=[3.75, 1.8, 1.9, 1.4, 1.4, 1.4], row_h=0.42)
    subtitle(s, "실험 조건의 통제", top=3.32)
    table(s, [
        ["조건", "처리", "확인 방법"],
        ["데이터", "동일 캐시 사용, API 재수집 없음", f"정규화 통계(mean/std) 일치 여부 검증: **{M.get('same_dataset')}**"],
        ["분할", "SEED 42 고정으로 동일한 학습/검증 분할", "random_split 생성기 시드 고정"],
        ["종료 기준", "동일한 조기 종료 기준(patience 20)", "양 조건 모두 최고 검증 시점의 가중치를 저장"],
        ["축 차단 방식", "아키텍처와 용량은 유지하고 Re·Im 입력만 영벡터로 대체(인코더 동결)",
         "게이트가 Z축에 **99.72%**를 배분하고 축 간 코사인이 모두 0.0000으로 측정됨"],
    ], BODY_L, 3.72, F, col_w=[1.5, 5.4, 4.75], row_h=0.4)
    _, tf = _txbox(s, BODY_L, 5.5, F, 0.8)
    p = tf.paragraphs[0]
    p.line_spacing = 1.2
    _run(p, "해석 — ", 10.5, bold=True, color=ACCENT)
    _run(p, "회귀 지표(기온·강수)의 개선 폭은 약 15% 수준인 반면, 희소 사건 분류에서 더 큰 개선이 관측되었다. "
            "황사 F1은 2.9배, 한파 F1은 1.7배로 증가하였다. 드문 사건일수록 단일 지점의 현재 상태만으로는 판별이 어려우며, "
            "주변 지역과 변화 경향 정보가 판별에 기여한 것으로 해석된다.\n", 10.5, color=BODY)
    _run(p, "한계 — ", 10.5, bold=True, color=WARN)
    _run(p, "조기 종료 기준은 동일하나 조건별 학습 길이는 상이할 수 있다(3축 78에폭, Z축 단독 25에폭). "
            "각 조건의 최고 검증 시점 가중치를 비교한 결과이다.", 10.5, color=BODY)
    kbox(s, [("동일 조건 통제 하에서 3축 구조가 회귀·분류 전 지표에서 우위를 보였으며, ", False),
             ("희소 사건 분류에서 그 효과가 가장 크게 나타났다", True), (".", False)])

    # ══════════════════════════════════ 11. 축 구성 변경 이력
    s = new_slide(prs, "2-10. 축 구성 변경 이력 — 입력 데이터 교체의 효과",
                  "아키텍처가 아닌 입력 데이터의 교체가 성능에 미친 영향", "적용 효과", 11)
    subtitle(s, "단계별 측정값 (v3~v7, 동일 검증셋 기준)")
    table(s, [
        ["단계", "Re축 입력", "Im축 입력", "기온 MAE", "강수 MAE", "황사 F1", "게이트 Re/Im/Z"],
        ["v3", "합성 위성 이미지", "MiniLM 임베딩", "1.450", "0.2175", "0.258", ".05 / .14 / .81"],
        ["v4", "합성 위성 이미지", "MiniLM 임베딩", "1.387", "0.2065", "0.257", "**.00** / .12 / .88"],
        ["**v5**", "**실측 IDW 보간장**", "MiniLM 임베딩", "1.279", "0.1696", "**0.661**", "**.52** / .04 / .44"],
        ["v6", "실측 IDW 보간장", "시간 경향 벡터(32)", "1.275", "0.1723", "0.593", ".41 / .15 / .44"],
        ["v7", "실측 IDW 보간장", "시간 경향 벡터(128)", "1.270", "0.1735", "0.576", ".37 / .21 / .42"],
    ], BODY_L, 1.42, F, col_w=[0.8, 2.35, 2.25, 1.5, 1.6, 1.25, 1.9], row_h=0.34)
    _, tf = _txbox(s, BODY_L, 3.55, F, 0.42)
    p = tf.paragraphs[0]
    _run(p, "※ v3~v7은 데이터가 3.6년으로 확장되기 이전 데이터셋에서 측정한 값이며, "
            "정규화 통계 불일치를 확인하였다. 위 표 내부의 비교만 유효하다.", 9.5, color=WARN)
    subtitle(s, "관측된 사항", top=4.0)
    table(s, [
        ["구간", "변경 내용", "결과"],
        ["v4 → v5", "Re축을 합성 위성에서 **실측 IDW 공간 보간장**으로 교체",
         "황사 F1 0.257 → **0.661** (2.6배). 게이트 배분 0.00 → **0.52**. 아키텍처는 변경하지 않았다"],
        ["v5 → v7", "Im축을 MiniLM 임베딩에서 **실측 시간 경향 벡터**로 교체",
         "게이트 배분 0.04 → 0.21로 증가하였으나 4개 F1 지표는 모두 하락하였다. "
         "인코더 용량을 32→128로 확대하여도 회복되지 않아 용량 부족 가설은 기각되었다"],
    ], BODY_L, 4.4, F, col_w=[1.2, 4.3, 6.15], row_h=0.62)
    kbox(s, [("Im축 교체는 성능 저하를 수반하였으나, ", False),
             ("측정되지 않은 데이터를 입력으로 사용하지 않는다는 원칙을 성능보다 우선", True),
             ("하여 유지하고 그 손실을 명시하였다.", False)])

    # ══════════════════════════════════ 12. 트러블슈팅 ①
    s = new_slide(prs, "2-11. 트러블슈팅 ① — 시뮬레이션 입력의 정보량",
                  "기여도 0의 확인 과정과 실측 데이터로의 교체", "트러블슈팅", 12)
    subtitle(s, "문제 · 대응 · 결과")
    table(s, [
        ["단계", "내용"],
        ["문제", "초기 Re축(합성 위성)과 Im축(MiniLM 예보문)은 **Z축 수치로부터 결정론적으로 생성**된 시뮬레이션이었다. "
                 "정보이론적으로 Z축을 초과할 수 없는 구조이나, 외형상으로는 3축 멀티모달 구성으로 보인다"],
        ["대응", "추론 시 Re축 입력을 **영벡터로 대체**하고 출력을 비교하였다(deploy_ablation.py). "
                 "동시에 게이트 배분값의 변화를 관측하였다"],
        ["결과", "축을 차단하여도 출력이 **소수점 이하까지 동일**하였다(기여도 0). 게이트 배분 역시 **0.000**으로 측정되었다. "
                 "모델이 해당 축의 무용성을 이미 학습한 상태였음을 의미한다"],
    ], BODY_L, 1.42, W, col_w=[0.85, 6.6], row_h=0.75)
    subtitle(s, "재설계 원칙과 적용 결과", top=4.05, left=0.59, width=7.6)
    table(s, [
        ["축", "교체 전", "교체 후", "결과"],
        ["Re", "합성 위성 이미지", "11개 지점 실측 IDW 보간장", "게이트 0.00 → **0.52**, 황사 F1 0.257 → **0.661**"],
        ["Im", "MiniLM 시뮬레이션 예보문", "1·3·6h 전 대비 실측 변화량", "게이트 0.12 → 0.21, F1 지표는 하락"],
    ], BODY_L, 4.45, W, col_w=[0.6, 2.0, 2.25, 2.6], row_h=0.45)
    _, tf = _txbox(s, BODY_L, 5.85, W, 0.45)
    p = tf.paragraphs[0]
    p.line_spacing = 1.15
    _run(p, "부수 효과 — ", 10, bold=True, color=ACCENT)
    _run(p, "sentence-transformers(약 450MB) 의존성이 서빙 경로에서 제거되어 배포 메모리가 227.9MiB로 축소되었다.",
         10, color=BODY)
    img_placeholder(s, "⑥ 모델 구조 탭 — 3축 정의",
                    "각 축이 참조하는 정보와\n입력 구성을 설명하는 영역")
    kbox(s, [("재설계 기준은 ", False),
             ("단일 시점 관측에서 유도할 수 없는 정보만을 새 축에 배치", True),
             ("한다는 것이다.", False)])

    # ══════════════════════════════════ 13. 트러블슈팅 ②
    s = new_slide(prs, "2-12. 트러블슈팅 ② — 학습 및 배포 단계의 결함",
                  "예외를 발생시키지 않아 실측을 통해서만 확인 가능하였던 사례", "트러블슈팅", 13)
    subtitle(s, "학습 단계")
    table(s, [
        ["관측된 현상", "원인", "조치"],
        ["강수 헤드의 학습이 중단됨",
         "clamp(x, 0, 200)이 x<0 구간에서 gradient가 정확히 0이 되어, 전 표본이 음수로 수렴한 이후 회복 불가",
         "softplus로 대체하고 회귀 테스트로 고정(gradient 0 대 0.475)"],
        ["강수가 실질적으로 학습되지 않음",
         "기온 손실(분산 수십 규모)이 강수 손실(O(1))을 약 1000배 초과",
         "분산 정규화를 적용하여 손실 규모를 정렬"],
        ["존재하지 않는 개선이 관측됨",
         "기준선을 전체 데이터에서 산출하여 검증셋과 표본이 불일치",
         "검증셋과 동일 표본에서 산출하도록 수정"],
        ["지점 간 잘못된 표본 쌍 생성 가능",
         "관측소 경계에서 (t, t+6) 쌍이 다른 지점으로 연결될 수 있는 인덱싱 구조",
         "(지점, 시각) 키 직접 조회로 저장 순서 의존성 제거"],
        ["**CI가 항상 통과로 표시됨**",
         "테스트가 실패를 출력만 하고 종료 코드 0을 반환",
         "sys.exit(0 if n_pass == n_all else 1) 추가"],
        ["9시간 이전 값이 실시간으로 표시됨",
         "컨테이너 기본 시간대가 UTC이므로 datetime.now()가 KST보다 9시간 지연",
         "모든 시각 연산에 ZoneInfo('Asia/Seoul') 명시"],
    ], BODY_L, 1.42, F, col_w=[3.0, 5.2, 3.45], row_h=0.44)
    subtitle(s, "배포 단계", top=4.5)
    table(s, [
        ["관측된 현상", "원인", "조치"],
        ["배포본이 대체값을 표시",
         "Secrets의 TOML 값에 인용부호가 누락되어 파싱 실패. 화면은 대체값으로 정상 표시되어 식별이 어려웠다",
         "API 호출 카운터를 화면에 노출하여 호출 0건으로 원인을 특정"],
        ["**새 체크포인트가 반영되지 않음**",
         "st.cache_resource가 인수 없이 적용되어 캐시 키가 불변. 프로세스가 유지되면 이전 모델이 계속 사용된다",
         "체크포인트의 수정시각·크기를 캐시 키에 포함하여 자동 무효화"],
    ], BODY_L, 4.9, F, col_w=[3.0, 5.2, 3.45], row_h=0.5)
    kbox(s, [("8건 모두 예외를 발생시키지 않았으며, ", False),
             ("실측 비교를 통해서만 확인 가능", True), ("하였다.", False)])

    # ══════════════════════════════════ 14. 산출 로직
    s = new_slide(prs, "2-13. +6시간 산출 로직",
                  "공유 표현으로부터 6개 헤드가 5종의 출력을 산출하는 구조", "산출 로직", 14)
    subtitle(s, "기온 — 퍼시스턴스 잔차 방식")
    code_line(s, "T(t+6)  =  T(t)  +  head_temp(s)", BODY_L, 1.38, W, size=12)
    bullets(s, [
        (0, [("절대 기온을 직접 산출하지 않는다. ", True),
             ("각 축이 단위 정규화되고 제곱근 연산이 부호를 제거하므로, 결합 표현만으로 절대 기온을 복원하는 것은 구조적으로 불리하다.", False)]),
        (0, [("T(t+6) ≈ T(t)는 강한 물리적 기준선", True),
             ("이므로, 이를 물리적으로 제공하고 모델은 그 편차만 학습한다. 논문의 오프셋 보정 구조와 동일하다.", False)]),
    ], top=1.86, width=W, size=10.5, bottom=2.9)
    subtitle(s, "강수 — Hurdle 구조", top=2.9, left=0.59, width=7.6)
    code_line(s, "y_precip  =  σ(z₁)  ×  softplus(z₂)", BODY_L, 3.3, W, size=12)
    table(s, [
        ["요소", "역할", "선정 근거"],
        ["σ(z₁)", "강수 발생 확률 (0~1)", "무강수가 대부분인 분포에서 **정확한 0**을 연속적으로 표현하기 위함"],
        ["softplus(z₂)", "발생 시 강수량 (항상 0 초과)", "강수량은 음수가 될 수 없다. clamp는 gradient 소실 위험이 있어 배제"],
        ["후처리", "0.2mm 미만은 0으로 반올림", "미세 노이즈 억제 (PRECIP_CLIP_THRESH)"],
    ], BODY_L, 3.78, W, col_w=[1.15, 2.3, 4.0], row_h=0.44)
    _, tf = _txbox(s, BODY_L, 5.35, W, 0.95)
    p = tf.paragraphs[0]
    p.line_spacing = 1.15
    _run(p, "실시간 입력 조회 — ", 10, bold=True, color=ACCENT)
    _run(p, "출력 1건당 15회(대상 지점 1, 인접 지점 11, 과거 1·3·6시간 3). "
            "SUCCESS_LIVE 응답만 입력에 포함한다. 폴백 레코드를 사용할 경우 합성 기본값(20.0°C/1013hPa)이 "
            "실측으로 간주되어 보간장에 반영되기 때문이다. 결측은 추정으로 보완하지 않고 중립값"
            "(보간장 0.5, 경향 0)으로 두며, 해당 사실을 화면에 표시한다.", 10, color=BODY)
    img_placeholder(s, "⑦ 출력값 추이 탭",
                    "최근 72시간 실측 그래프와\n★ 표시된 +6시간 산출값")
    kbox(s, [("물리적 제약(음수 불가, 강한 기준선의 존재)을 ", False),
             ("손실 함수가 아닌 아키텍처 수준에서 반영", True), ("하였다.", False)])

    # ══════════════════════════════════ 15. 극한기상
    s = new_slide(prs, "2-14. 극한기상 3종 — 라벨 정의와 판정 임계값",
                  "정답 라벨의 출처와 임계값 재보정 절차", "산출 로직", 15)
    subtitle(s, "라벨 정의")
    table(s, [
        ["사건", "학습에 사용한 정답 라벨", "기록 부재 시 대체", "판정선", "F1"],
        ["폭염", "**기상청이 발표한 폭염주의보 기록**", "발표기준 수치 **33°C**를 대상 시각 기온에 적용", "0.84", f1("heatwave")],
        ["한파", "**발표된 한파주의보 기록**", "발표기준 수치 **−12°C**를 대상 시각 기온에 적용", "0.85", f1("coldwave")],
        ["황사", "**PM10 1시간 평균 150㎍/㎥ 이상** (대기환경지수 '매우나쁨')", "대체 없음. 마스크로 손실에서 제외", "0.91", f1("dust")],
    ], BODY_L, 1.42, W, col_w=[0.65, 3.0, 2.4, 0.7, 0.7], row_h=0.62)
    _, tf = _txbox(s, BODY_L, 3.95, W, 1.0)
    p = tf.paragraphs[0]
    p.line_spacing = 1.15
    _run(p, "라벨 정의상의 제약 — ", 10, bold=True, color=WARN)
    _run(p, "공식 황사주의보 기준(400㎍/㎥ 2시간 지속)은 보유 데이터의 6개 지점에서 사례가 거의 없어 학습이 불가능하였다. "
            "150㎍/㎥ 기준은 815건(0.56%)으로 한파와 유사한 규모이다. 또한 33°C·−12°C는 본래 일 최고기온·아침 최저기온에 "
            "적용되는 기준이나 대체 라벨에서는 시간별 기온에 적용하므로, 발표된 특보와 동일한 정의가 아니다. "
            "두 사항 모두 화면에 명시한다.", 10, color=BODY)
    subtitle(s, "임계값 재보정 절차", top=5.0, left=0.59, width=7.6)
    bullets(s, [
        (0, [("① 검증셋을 보정용·평가용으로 50:50 무작위 분할(학습/검증 분할과 별개의 시드 사용)", False)]),
        (0, [("② 보정용에서만 F1을 최대화하는 임계값을 탐색", False)]),
        (0, [("③ 미사용 평가용에 적용하여 성능 유지 여부를 확인. 미달 시 과적합으로 판정하여 폐기", False)]),
    ], top=5.4, width=W, size=10, bottom=6.3)
    img_placeholder(s, "⑧ 극한 기상 탭",
                    "폭염·한파·황사 확률 막대와\n판정 임계값 표시선")
    kbox(s, [("강수 임계값 재보정은 이 검증을 통과하지 못하여(순이득 0.01 미만) ", False),
             ("기존 값을 유지", True), ("하였으며, 폭염·한파·황사만 채택하였다.", False)])

    # ══════════════════════════════════ 16. 성능 검증
    s = new_slide(prs, "2-15. 성능 검증",
                  "기준선 대비 달성 지표와 미달 지표", "성능 검증", 16)
    subtitle(s, "회귀 성능 — 검증셋 66,319건, 학습에 미사용")
    table(s, [
        ["지표", "모델", "기준선 (학습을 수행하지 않는 자명한 방법)", "판정"],
        ["기온 MAE", f"**{M['tri_temp']:.4f} °C**",
         f"퍼시스턴스 T(t+6)≈T(t) : {fnum(M['naive_temp'], '{:.4f}')} °C",
         f"**{pct(M['tri_temp'], M['naive_temp'])}**"],
        ["강수 MAE", f"**{M['tri_precip']:.4f} mm**",
         f"상시 0mm 산출 : {fnum(M['naive_precip'], '{:.4f}')} mm",
         f"**{pct(M['tri_precip'], M['naive_precip'])}**"],
    ], BODY_L, 1.42, W, col_w=[1.1, 1.5, 3.5, 1.35], row_h=0.44)
    _, tf = _txbox(s, BODY_L, 2.9, W, 1.15)
    p = tf.paragraphs[0]
    p.line_spacing = 1.15
    _run(p, "강수 지표에 대한 해석 — ", 10, bold=True, color=WARN)
    _run(p, "관측 시점의 대부분이 무강수이므로 '상시 0mm' 기준선의 MAE도 낮게 산출된다. 현행 모델은 이 기준선을 "
            "근소하게 상회하지 못하였다. 원인은 극한기상 헤드가 3종으로 증가하면서 헤드 간 공유 표현 경쟁이 발생하여 "
            "Hurdle 구조의 무강수 억제력이 약화된 것으로 확인되었다(초기 −25.6%). 손실 가중치 조정(1.0→0.3)과 "
            "데이터 확장(3.6년)을 통해 현재 수준까지 축소하였다. 강수는 MAE보다 발생 여부 판정으로 평가하는 것이 "
            "실용적이며, 해당 지표는 출력값 적중률의 강수 항목에 해당한다.", 10, color=BODY)
    subtitle(s, "극한기상 분류 성능 — 판정 임계값 0.5 기준", top=4.15, left=0.59, width=7.6)
    table(s, [
        ["사건", "정밀도", "재현율", "F1", "검증셋 발생"],
        ["폭염", pr("heatwave"), rc("heatwave"), f1("heatwave"), npos("heatwave")],
        ["한파", pr("coldwave"), rc("coldwave"), f1("coldwave"), npos("coldwave")],
        ["황사", pr("dust"), rc("dust"), f1("dust"), npos("dust")],
    ], BODY_L, 4.55, W, col_w=[1.0, 1.55, 1.55, 1.5, 1.85], row_h=0.32)
    _, tf = _txbox(s, BODY_L, 5.95, W, 0.42)
    p = tf.paragraphs[0]
    p.line_spacing = 1.1
    _run(p, "F1의 상한은 표본 수에 의해 제약된다. 황사의 값이 낮은 것은 학습 가능한 사례 수가 적기 때문이다. "
            "출력값 적중률은 기온 ±1.5°C, 강수 0.1mm 기준 발생 여부 일치로 판정한다.", 9.5, color=SUB)
    img_placeholder(s, "⑨ 성능 검증 탭",
                    "회귀 MAE와 기준선 대비 증감,\n극한기상 F1 및 적중률")
    kbox(s, [("기준선 수치를 체크포인트에 저장하도록 개선하여, ", False),
             ("지표의 상대적 위치가 화면에 자동 표시", True), ("된다.", False)])

    # ══════════════════════════════════ 17. 활용 시나리오
    s = new_slide(prs, "2-16. 활용 시나리오 — 산출값과 의사결정의 연계",
                  "관측 지점이 위치한 12개 도시에서의 적용 가능 범위와 적용 한계", "활용", 17)
    subtitle(s, "산출값별 활용 — 12개 지점은 인구·산업이 밀집한 지역으로, 기후 리스크에 노출된 자산이 존재한다")
    table(s, [
        ["산출값", "제공하는 정보", "활용 예"],
        ["**기온** (현재·+6h)", "6시간 뒤 기온 수준. 기준선 대비 유의한 개선이 확인된 회귀 지표",
         "야외 작업의 온열질환 예방(작업 중지·휴식 주기 사전 조정), 태양광 발전 효율 예측(패널 온도 상승 시 효율 저하), 냉난방 피크 수요 대비"],
        ["**강수** (발생 여부)", "강수 발생 여부. 강수량 절대값이 아닌 이진 판정으로 활용한다",
         "야외 작업 및 콘크리트 타설 일정 조정, 발전량 급감 대비, 배수 설비 사전 점검"],
        ["**습도**", "체감온도(열지수) 산정의 필수 구성 요소",
         "기온 단독으로는 온열질환 위험을 판정할 수 없다. 습도가 높을 경우 땀 증발이 저해되어 동일 기온에서도 위험이 상승한다"],
        ["**기압**", "기압 하강률은 저기압 접근의 선행 지표에 해당한다 (Im축이 사용하는 신호)",
         "급변 감시. 절대값보다 변화 방향이 정보를 제공한다"],
        ["**폭염 확률**", "폭염특보 수준 상황의 6시간 뒤 발생 확률",
         "작업 중지 판단, 전력 피크 수요 대비, 온열질환 감시 대상 선별"],
        ["**한파 확률**", "한파특보 수준 상황의 발생 확률",
         "동파·결빙 대비, 난방 수요 예측, 취약계층 보호 자원 배치"],
        ["**황사 확률**", "PM10 '매우나쁨'(150㎍/㎥) 수준의 발생 확률",
         "야외 작업자 보호구 지급, 태양광 패널 오염에 따른 세척 시점 판단, 실내 공기질 설비 가동"],
    ], BODY_L, 1.42, F, col_w=[1.55, 3.7, 6.4], row_h=0.5)
    _, tf = _txbox(s, BODY_L, 5.05, F, 1.3)
    p = tf.paragraphs[0]
    p.line_spacing = 1.15
    _run(p, "적용 한계 — 활용 범위와 함께 명시한다.  ", 10, bold=True, color=WARN)
    _run(p, "① 강수량 절대값에 근거한 의사결정에는 사용할 수 없다. 기준선을 상회하지 못하였으므로 발생 여부 판정으로만 활용한다.  "
            "② 12개 지점 외 지역에는 산출값이 존재하지 않는다.  "
            "③ 황사 라벨은 6개 지점(서울·수원·춘천·대구·전주·광주)에만 존재하며, 나머지 지점의 확률은 타 지역에서 학습된 "
            "패턴을 적용한 것이므로 신뢰도가 낮다.  "
            "④ 공식 예보·특보를 대체하지 않는다. 기상산업진흥법상 예보업은 등록 대상이며, 본 산출물은 모델 출력값이다.\n",
         10, color=BODY)
    _run(p, "지역 특성의 반영 — ", 10, bold=True, color=ACCENT)
    _run(p, "강릉(동해 영향), 제주(해양성), 대구(분지 고온) 등 동일한 기온이라도 지역별 의미가 상이하다. "
            "모델은 위도·경도를 입력으로 포함하여 이 차이를 학습한다.", 10, color=BODY)
    kbox(s, [("기후 리스크 정보는 현재 상태의 제시에 그치지 않고 ", False),
             ("일정 시간 이후의 대응 판단으로 연결될 때 실무적 가치를 갖는다", True), (".", False)])

    # ══════════════════════════════════ 18. 배포·운영
    s = new_slide(prs, "2-17. 배포 및 운영 — API 호출량 설계와 갱신 자동화",
                  "캐시 키를 관측 시각에 연동하여 호출량을 접속자 수와 무관하게 설계", "운영", 18)
    subtitle(s, "API 호출량 — 출력 1건당 15회 (대상 1 + 인접 11 + 과거 3)")
    table(s, [
        ["조건", "시간당", "일간", "평가"],
        ["캐시 미적용, 5분 갱신, 1개 지점", "180", "4,320", "제한 근접"],
        ["캐시 미적용, 5분 갱신, 접속자 3인", "540", "12,960", "**차단 지점(9,800) 초과**"],
        ["**캐시 적용** (지점·관측시각 키)", "**15**", "**360**", "차단 지점의 3.7%"],
    ], BODY_L, 1.42, F, col_w=[4.0, 1.3, 1.5, 4.85], row_h=0.32)
    bullets(s, [
        (0, [("설계 근거는 관측 주기라는 물리적 사실이다. ", True),
             ("ASOS는 매시 정각에 1회 관측하며 약 10분 후 공개된다. 동일 정시 내에서는 재조회하여도 값이 변하지 않으므로, "
              "캐시 키를 관측 시각으로 두면 호출량이 갱신 주기 및 접속자 수와 무관해진다.", False)]),
        (0, [("검증 결과: 동일 정시에 2회 실행 시 호출 증가 0건. ", True),
             ("2차 안전장치로 일간 상한(기본 3,000회)을 두고 사용량을 화면에 표시한다. "
              "합산 예산은 대시보드 360~720회와 자동 갱신 288회를 포함하여 약 650~1,000회/일이며, 차단 지점의 7~10% 수준이다.", False)]),
    ], top=3.0, size=10.5, bottom=4.2)
    subtitle(s, "재학습 자동화 검토", top=4.2)
    table(s, [
        ["단계", "자동화 가능성", "근거"],
        ["신규 관측 수집", "**적용 완료**", "GitHub Actions 6시간 주기 실행 후 저장소 커밋, 배포 환경 자동 갱신"],
        ["극한기상 라벨 갱신", "부분 자동화", "날씨 이슈 데이터 및 PM10은 **API가 제공되지 않아** 포털에서 수동 수집이 필요하다"],
        ["재학습 실행", "로컬 GPU 필요", "무료 러너에 GPU가 없다. **권장 주기는 분기~반기** — 월 단위는 신규 표본이 전체의 약 2%로 이득이 제한적이다"],
        ["검증·배포 게이트", "적용 가능", "재학습 결과가 기존 지표에 미달할 경우 배포를 차단하는 CI 게이트(설계 완료)"],
    ], BODY_L, 4.58, F, col_w=[1.85, 1.5, 8.3], row_h=0.34)
    kbox(s, [("인증 정보는 .env(로컬), Streamlit Secrets(배포), GitHub Secrets(자동화)로 ", False),
             ("세 경로 분리", True), ("되어 저장소에 포함되지 않으며, API 호출은 서버 측에서만 수행된다.", False)])

    # ══════════════════════════════════ 19. 한계 및 향후 과제
    s = new_slide(prs, "2-18. 한계 및 향후 과제",
                  "현 시점에서 확인된 제약과 후속 연구 방향", "한계", 19)
    subtitle(s, "확인된 한계")
    table(s, [
        ["항목", "내용", "현 상태"],
        ["**강수 지표의 기준선 미달**", "'상시 0mm' 기준선을 상회하지 못하였다. 헤드 간 표현 경쟁이 원인으로 확인되었다",
         "손실 가중치 조정과 데이터 확장으로 격차를 축소하였으나 해소되지 않았다"],
        ["**Im축 교체에 따른 성능 저하**", "실측 경향 벡터로 교체한 이후 4개 F1 지표가 모두 하락하였다. 인코더 용량 확대로 회복되지 않았다",
         "원칙(측정되지 않은 데이터 배제)을 우선하여 유지하고 손실을 명시하였다"],
        ["**Gram-Schmidt 직교화의 부적합**", "원 논문은 축 독립성을 전제하나, 본 도메인에서는 직교화 강제 시 성능이 하락하였다",
         "기본값을 ORTHOGONALIZE=False로 설정"],
        ["**위상각의 수학적 퇴화**", "단위 정규화 구현에서 θ는 게이트 가중치의 재표현으로 축소된다(차이 1.19e-07)",
         "상태 표현에서 축 간 코사인 유사도로 대체"],
        ["**PPO의 기준선 미달**", "경보 게이트에서 PPO(−0.173)가 지도학습(−0.112) 및 고정 임계값(−0.139)에 미달하였다",
         "이벤트 214건으로 표본이 부족한 것으로 분석"],
    ], BODY_L, 1.42, F, col_w=[2.4, 5.0, 4.25], row_h=0.56)
    subtitle(s, "향후 과제", top=4.22)
    table(s, [
        ["과제", "내용"],
        ["실측 위성 자료 연계", "천리안 GK2A 또는 Sentinel-2를 Re축에 적용하여 12개 지점 보간의 공간 해상도 제약을 개선한다"],
        ["레이더 자료 기반 초단기 예측", "대류성 강수는 지상 관측만으로 예측 가능 시계가 제한적이다. 반사도 자료를 도입하여 강수 성능을 개선한다"],
        ["이상 탐지 계층 도입", "현재는 결측 레코드를 제외하는 방식이다. Autoencoder 복원오차 기반 센서 이상 탐지를 전처리 단계에 배치한다"],
        ["물리 제약 학습(PINN)", "열역학·질량보존을 손실 제약으로 도입한다. 다만 보간장에 이류·확산을 적용하면 자체 인공물을 학습하므로 실측 격자 확보가 선행되어야 한다"],
        ["재학습 게이트 자동화", "검증 지표가 기존 대비 미달할 경우 배포를 차단하는 CI 게이트. 설계는 완료되었으며 구현이 남아 있다"],
    ], BODY_L, 4.60, F, col_w=[2.4, 9.25], row_h=0.28)
    kbox(s, [("본 연구는 성능 달성보다 ", False),
             ("설계 결함을 실험으로 규명하고 그 근거를 제시하는 과정", True), ("에 중점을 두었다.", False)])

    prs.save(OUT)
    print(f"생성 완료: {OUT}  (슬라이드 {len(prs.slides)}장)")
    print(f"  3축   기온 {M['tri_temp']:.4f} / 강수 {M['tri_precip']:.4f}")
    print(f"  Z단독 기온 {fnum(M['z_temp'], '{:.4f}')} / 강수 {fnum(M['z_precip'], '{:.4f}')}")
    print(f"  기준선 기온 {fnum(M['naive_temp'], '{:.4f}')} / 강수 {fnum(M['naive_precip'], '{:.4f}')}")
    print(f"  동일 데이터셋 여부: {M.get('same_dataset')}")


if __name__ == "__main__":
    build()
